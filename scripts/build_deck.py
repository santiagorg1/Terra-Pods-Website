"""ISN deck v3 builder — surgical edits on slides 7/9/10/11/12/17.

Loads the existing ISN-Customs-Brokerage-Deck.pptx as a template, then:
  - swaps the picture on slides 7, 9, 10, 11, 12 (cover-cropped to original placement)
  - replaces the slide-17 text-card grid with a 5x3 client-logo grid
Everything else (eyebrows, headlines, bullets, overlays, footers) is left untouched.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
PPTX = os.path.join(REPO, "public/presentations/ISN-Customs-Brokerage-Deck.pptx")

PHOTO_SWAPS = [
    (7,  "public/isn/photos/08-customs-officer.jpg"),
    (9,  "public/isn/photos/10-aerial-logistics.jpg"),
    (10, "public/isn/photos/04-air-cargo.jpg"),
    (11, "public/isn/photos/03-container-ship.jpg"),
    (12, "public/isn/photos/02-highway-truck.jpg"),
]

SLIDE_17_LOGOS = [
    "current.png", "terrapods.png", "resco.png", "commercial-solutions.png", "nafta.png",
    "underterra.png", "becksteel.png", "sanford.png", "doerfer.png", "ge.png",
    "parkohio.png", "csi.png", "general-aluminum.png",
]
ROW_COUNTS = [5, 5, 3]  # 5+5+3 = 13


def cover_crop(src_w, src_h, tgt_w, tgt_h):
    """python-pptx crop fractions (left, top, right, bottom) for cover-scale."""
    src_r = src_w / src_h
    tgt_r = tgt_w / tgt_h
    if src_r < tgt_r:
        c = (1 - src_r / tgt_r) / 2
        return 0.0, c, 0.0, c
    c = (1 - tgt_r / src_r) / 2
    return c, 0.0, c, 0.0


def replace_picture(slide, new_path):
    """Find first PICTURE on slide; replace with new image at same placement + z-order, cover-cropped."""
    pic = next((s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
    if pic is None:
        raise RuntimeError("no picture shape on slide")
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    sp_tree = slide.shapes._spTree
    z_idx = list(sp_tree).index(pic._element)
    sp_tree.remove(pic._element)
    new_pic = slide.shapes.add_picture(new_path, left, top, width, height)
    # Move to original z-order position
    el = new_pic._element
    sp_tree.remove(el)
    sp_tree.insert(z_idx, el)
    # Cover-crop the new image to placement aspect
    with Image.open(new_path) as im:
        sw, sh = im.size
    cl, ct, cr, cb = cover_crop(sw, sh, width, height)
    new_pic.crop_left = cl
    new_pic.crop_top = ct
    new_pic.crop_right = cr
    new_pic.crop_bottom = cb


def is_client_card(shp):
    """True for shapes that compose the old slide-17 client text-card grid."""
    if shp.top is None or shp.height is None or shp.width is None:
        return False
    top_in = shp.top / 914400
    w_in = shp.width / 914400
    h_in = shp.height / 914400
    if not (3.4 <= top_in <= 6.5):
        return False
    if not (0.6 <= h_in <= 0.8):
        return False
    return 3.6 <= w_in <= 4.1


def rebuild_slide_17(slide):
    # Remove old card shapes
    to_remove = [s for s in slide.shapes if is_client_card(s)]
    for s in to_remove:
        s._element.getparent().remove(s._element)

    # Logo grid geometry
    slide_w = 13.33
    grid_top = 3.45
    grid_bottom = 6.85
    grid_h = grid_bottom - grid_top
    cell_h = grid_h / len(ROW_COUNTS)
    col_count = max(ROW_COUNTS)
    side_margin = 0.60
    usable_w = slide_w - 2 * side_margin
    cell_w = usable_w / col_count
    pad_x = 0.18
    pad_y = 0.13
    max_w = cell_w - 2 * pad_x
    max_h = cell_h - 2 * pad_y

    idx = 0
    for row, count in enumerate(ROW_COUNTS):
        row_w = count * cell_w
        row_left = (slide_w - row_w) / 2
        for col in range(count):
            if idx >= len(SLIDE_17_LOGOS):
                return
            logo_path = os.path.join(REPO, "public/isn/clients", SLIDE_17_LOGOS[idx])
            with Image.open(logo_path) as im:
                iw, ih = im.size
            img_r = iw / ih
            box_r = max_w / max_h
            if img_r >= box_r:
                w = max_w
                h = max_w / img_r
            else:
                h = max_h
                w = max_h * img_r
            cell_left = row_left + col * cell_w
            cell_top = grid_top + row * cell_h
            x = cell_left + (cell_w - w) / 2
            y = cell_top + (cell_h - h) / 2
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y), Inches(w), Inches(h))
            idx += 1


def main():
    prs = Presentation(PPTX)
    for slide_no, rel_path in PHOTO_SWAPS:
        path = os.path.join(REPO, rel_path)
        assert os.path.exists(path), f"missing photo: {path}"
        replace_picture(prs.slides[slide_no - 1], path)
        print(f"slide {slide_no}: bg -> {rel_path}")
    rebuild_slide_17(prs.slides[16])
    print(f"slide 17: rebuilt with {len(SLIDE_17_LOGOS)} client logos in {ROW_COUNTS} grid")
    prs.save(PPTX)
    print(f"\nsaved {PPTX}")


if __name__ == "__main__":
    main()
